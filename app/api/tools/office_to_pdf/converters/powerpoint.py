from __future__ import annotations

from io import BytesIO

import fitz
from pptx import Presentation
from pptx.util import Pt


def convert_to_powerpoint(pdf_path: str, output_path: str) -> None:
    prs = Presentation()
    if len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    doc = fitz.open(pdf_path)
    try:
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            rect = page.rect

            prs.slide_width = Pt(rect.width)
            prs.slide_height = Pt(rect.height)

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            zoom = 2
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)

            img_bytes = pix.tobytes("png")
            image_stream = BytesIO(img_bytes)

            slide.shapes.add_picture(
                image_stream,
                0,
                0,
                width=Pt(rect.width),
                height=Pt(rect.height),
            )
    finally:
        doc.close()

    prs.save(output_path)
